from __future__ import annotations

from pathlib import Path
from pptx import Presentation


ROOT = Path('/Users/timothycllam/Documents/llm_sandbox/openai/codex_assignment/retailnext_outfit_assistant')
TEMPLATE = Path('/Users/timothycllam/Downloads/AI-Powered_Product_Discovery_for_RetailNext.pptx')
OUTPUT = ROOT / 'deliverables' / 'AI-Powered_Product_Discovery_for_RetailNext_FINAL.pptx'

SCREENSHOT_MAIN = ROOT / 'assets' / 'screenshots' / '05_image_match_results.png'
SCREENSHOT_MODAL = ROOT / 'assets' / 'screenshots' / '04_upload_modal.png'


SLIDE_TEXT = {
    1: {
        2: 'RetailNext Outfit Assistant: AI-Powered Product Discovery',
        3: 'Executive + CTO briefing for reducing product-findability complaints using multimodal search and match intelligence.',
    },
    2: {
        2: 'Business Situation at RetailNext',
        5: 'Problem 1',
        6: 'Findability',
        7: 'Recent customer reviews cite difficulty finding updated styles or specific items before upcoming events.',
        10: 'Problem 2',
        11: 'Intent Gap',
        12: 'Shoppers describe looks in natural language or photos; keyword/filter search misses this intent.',
        15: 'Problem 3',
        16: 'Store Burden',
        17: 'Associates spend time manually translating vague requests into inventory searches, slowing conversion.',
    },
    3: {
        2: 'What We Built (Demo Ready Today)',
        3: 'A full-stack prototype built from the OpenAI cookbook sample_clothes dataset and extended for RetailNext.',
        4: 'Natural-Language Query Search',
        5: "Bob types intent (for example: 'my wife wants a sakura season t shirt') and receives ranked similar items.",
        6: 'Image Upload Match',
        7: 'Customer uploads a photo; vision analysis extracts style attributes and generates retrieval-ready search queries.',
        8: 'AI Match Explainability (Added Feature)',
        9: "'Check Your Match' returns verdict, rationale, and confidence for each item and persists results for follow-up.",
        11: '',
        13: '',
        15: '',
        16: '',
    },
    4: {
        2: 'OpenAI Platform Used in This Solution',
        3: 'Multimodal Understanding',
        4: 'Model: gpt-4o-mini',
        5: 'Analyzes uploaded outfit images and returns structured JSON (gender, occasion, colors, article types, search queries).',
        6: 'Key Capability: Vision + Structured Output',
        7: 'Semantic Retrieval',
        8: 'Model: text-embedding-3-large',
        9: 'Embeds shopper text and generated queries; cosine similarity ranks best matches across 1,000 catalog vectors.',
        10: 'Key Capability: Intent-Based Discovery',
        11: 'Session-Aware Match Scoring',
        12: 'Model/API: gpt-4o-mini via Responses API',
        13: "Evaluates each recommendation against session intent for 'Check Your Match' verdict + rationale + confidence.",
        14: 'Key Capability: Explainable Recommendations',
    },
    5: {
        2: 'Solution Architecture',
        3: 'Touchpoints',
        9: 'Web Storefront',
        10: 'Home + Personalized UI',
        16: 'Associate/Kiosk Ready',
        17: 'Same API-backed flows',
        18: 'Orchestration',
        20: 'FastAPI Service Layer',
        21: 'Python / FastAPI',
        22: '• Search + image-match APIs\n• Session context + ranking\n• Match-check persistence',
        23: 'Intelligence & Data',
        25: 'gpt-4o-mini',
        26: 'Vision extraction + match reasoning',
        28: 'text-embedding-3-large',
        29: 'Query + catalog embeddings',
        31: 'SQLite + Vector Cache',
        32: 'catalog_index.npz + session DB',
    },
    6: {
        3: 'Process Flow: Query to Recommendation',
        6: 'User Intent',
        7: 'Customer types a natural-language request or uploads an outfit image.',
        8: 'Web UI',
        11: 'Intent Extraction',
        12: 'Text path embeds query. Image path uses gpt-4o-mini to create structured search queries.',
        13: 'OpenAI API',
        16: 'Retrieval',
        17: 'Service runs cosine similarity on catalog embeddings and selects top candidates.',
        18: 'Numpy Vector Index',
        21: 'Ranking + Explain',
        22: "Product candidates are ranked; optional 'Check Your Match' call returns verdict/rationale/confidence.",
        23: 'Session Reasoner',
        26: 'Response + Logging',
        27: 'Personalized results render in UI; sessions and match checks are stored in SQLite for analytics.',
        28: 'UI + Persistence',
    },
    7: {
        2: 'Business Value & Pilot KPIs',
        3: '15-20%',
        4: 'Conversion Lift Target',
        5: 'Expected improvement from better relevance on event-driven and long-tail style requests.',
        6: '30-45%',
        7: 'Faster Time-to-Item',
        8: 'Reduction in search friction using natural language and image-driven discovery.',
        9: '10-15%',
        10: 'Return Rate Reduction',
        11: 'Match explanation can improve confidence before checkout and reduce mismatched purchases.',
        14: 'Strategic Impact',
        15: 'Capture structured customer intent data (occasion, color, style) for merchandising decisions.',
        16: "Increase loyalty by turning 'can't find it' moments into guided shopping experiences.",
        17: 'Create a reusable AI foundation across ecommerce, app, and in-store channels.',
        18: 'Operational Gains',
        19: 'Reduce associate lookup time and escalation for ambiguous customer requests.',
        20: 'Track query-to-click and match outcomes for continuous ranking improvements.',
        21: 'Roll out safely with API controls, monitoring, and deterministic fallback logic.',
    },
    8: {
        2: 'Implementation Roadmap',
        3: 'De-risked path from pilot to enterprise scale in 12 weeks.',
        6: 'Phase 1',
        7: 'Foundation',
        8: 'Weeks 1-2',
        10: 'Catalog sync + embedding cache build',
        12: 'API + session schema hardening',
        14: 'Security and guardrail baseline',
        16: 'Success metrics definition',
        19: 'Phase 2',
        20: 'Pilot',
        21: 'Weeks 3-6',
        23: 'Launch in 5 high-traffic stores',
        25: 'Associate enablement and playbooks',
        27: 'Weekly prompt/ranking tune-ups',
        29: 'Collect user feedback + query logs',
        32: 'Phase 3',
        33: 'Optimize',
        34: 'Weeks 7-9',
        36: 'Latency and caching improvements',
        38: 'Ranking quality improvements',
        40: 'KPI dashboard for funnel + match',
        42: 'A/B test against current search',
        45: 'Phase 4',
        46: 'Scale',
        47: 'Week 10+',
        49: 'Regional rollout expansion',
        51: 'Inventory/availability integrations',
        53: 'Seasonal campaign activation',
        55: 'Production governance cadence',
    },
    9: {
        2: 'Why OpenAI Platform for RetailNext',
        3: 'One platform for multimodal reasoning, embeddings, and response generation.',
        4: 'This enabled a fast build from cookbook dataset to full-stack prototype: understand intent, retrieve products, explain matches, and log outcomes.',
        5: 'Model Breadth',
        6: 'Vision, text generation, and embeddings are available in a consistent API surface for faster iteration.',
        7: 'Enterprise Controls',
        8: 'API-level controls and governance patterns support secure rollout for retail customer experiences.',
        9: 'Developer Velocity',
        10: 'SDKs and cookbook assets reduced implementation time while keeping architecture extensible.',
    },
    10: {
        2: 'Next Steps',
        3: 'Submission-ready plan plus onsite presentation flow.',
        4: 'Immediate Actions',
        6: 'Finalize Submission Package',
        7: 'Record a <6 minute technical walkthrough covering architecture, key code paths, and live UX flows.',
        9: 'Prepare Executive Demo',
        10: 'Use a 20-minute narrative: business problem, OpenAI platform, architecture, value, then live demo.',
        12: 'Agree Pilot Success Criteria',
        13: 'Align KPI targets, risk controls, and go/no-go thresholds with Head of Innovation and CTO.',
        14: 'Presentation Roles',
        15: 'Lead Presenter',
        16: '[Your Name]',
        17: 'Business + technical narrative',
        18: 'Demo/Q&A Lead',
        19: '[Your Name]',
        20: 'Architecture, code, and risk questions',
    },
}


SLIDE_NOTES = {
    1: (
        'Open with the business framing: RetailNext has a findability problem, not a product assortment problem. '
        'Position the solution as a practical way to convert intent into purchasable items in seconds.'
    ),
    2: (
        'State assumptions from discovery: poor reviews are concentrated in event-driven shopping journeys and long-tail style requests. '
        'Current keyword search and manual associate lookup create friction and missed conversions.'
    ),
    3: (
        'Walk through what was built: text search, image match, and the added Check Your Match feature. '
        'Use the screenshot as proof that the experience is functional end-to-end today.'
    ),
    4: (
        'Explain platform mapping to business function: vision for intent extraction, embeddings for retrieval, and response reasoning for explanation. '
        'Highlight that each model call has a specific job in the pipeline.'
    ),
    5: (
        'Describe architecture left to right: touchpoint, API orchestration, AI/data services. '
        'Mention persistence in SQLite for sessions and match checks to support analytics and future optimization.'
    ),
    6: (
        'Narrate the exact request lifecycle. Emphasize that image and text converge into the same semantic retrieval path, '
        'which keeps implementation simple and scalable.'
    ),
    7: (
        'Present these as pilot KPI targets, not guaranteed outcomes. '
        'Tie each KPI to a measurable mechanism in the built product (relevance, speed, confidence).' 
    ),
    8: (
        'Show risk-managed execution: foundation, pilot, optimization, scale. '
        'Point out feedback loops and A/B testing before broad rollout.'
    ),
    9: (
        'Frame OpenAI as the enabling platform: unified APIs, enterprise control patterns, and rapid developer iteration. '
        'Avoid abstract hype; tie claims back to what was implemented in this prototype.'
    ),
    10: (
        'Close with concrete next steps and role clarity for the onsite format. '
        'Invite questions on architecture decisions, metrics assumptions, and rollout risks.'
    ),
}


def set_shape_text(slide, shape_index: int, value: str) -> None:
    shape = slide.shapes[shape_index]
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = value


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f'Missing template deck: {TEMPLATE}')
    if not SCREENSHOT_MAIN.exists() or not SCREENSHOT_MODAL.exists():
        raise FileNotFoundError('Missing one or more captured screenshots.')

    prs = Presentation(str(TEMPLATE))

    for slide_no, mapping in SLIDE_TEXT.items():
        slide = prs.slides[slide_no - 1]
        for shape_index, text in mapping.items():
            set_shape_text(slide, shape_index, text)

    # Add real app screenshots to slide 3 (solution overview)
    slide3 = prs.slides[2]
    slide3.shapes.add_picture(
        str(SCREENSHOT_MAIN),
        left=4_880_000,
        top=1_120_000,
        width=3_520_000,
        height=2_180_000,
    )
    slide3.shapes.add_picture(
        str(SCREENSHOT_MODAL),
        left=5_250_000,
        top=3_360_000,
        width=2_760_000,
        height=1_520_000,
    )

    # Add speaker notes
    for slide_no, note_text in SLIDE_NOTES.items():
        slide = prs.slides[slide_no - 1]
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = note_text

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f'Wrote: {OUTPUT}')


if __name__ == '__main__':
    main()
